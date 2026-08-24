#!/usr/bin/env python3
"""Build the 24 August 2026 advisor update PowerPoint."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "slides" / "advisor-update-2026-08-24.pptx"

NAVY = RGBColor(23, 54, 93)
RED = RGBColor(155, 44, 44)
GREEN = RGBColor(33, 110, 57)
DARK = RGBColor(23, 32, 51)
MUTED = RGBColor(75, 85, 99)
LIGHT = RGBColor(247, 248, 250)
PALE_BLUE = RGBColor(223, 231, 242)
WHITE = RGBColor(255, 255, 255)


def set_background(slide, color=LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.35), Inches(12), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(29)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        sub = slide.shapes.add_textbox(
            Inches(0.68), Inches(1.08), Inches(11.9), Inches(0.45)
        )
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = "Aptos"
        p.font.size = Pt(15)
        p.font.color.rgb = MUTED


def add_footer(slide, number):
    box = slide.shapes.add_textbox(
        Inches(0.68), Inches(7.12), Inches(11.9), Inches(0.22)
    )
    p = box.text_frame.paragraphs[0]
    p.text = f"Faiza · CDC-aware RTL generation · 24 Aug 2026                                      {number}"
    p.font.name = "Aptos"
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED


def add_bullets(slide, bullets, x=0.85, y=1.55, w=11.6, h=4.9, size=22):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.name = "Aptos"
        p.font.size = Pt(size - 2 * level)
        p.font.color.rgb = DARK
        p.space_after = Pt(10 if level == 0 else 4)
        p.line_spacing = 1.05
    return box


def add_callout(slide, text, y=5.65, color=RED):
    shape = slide.shapes.add_shape(
        1, Inches(0.85), Inches(y), Inches(11.65), Inches(0.82)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALE_BLUE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = color


def add_table(slide, rows, widths, x=0.7, y=1.55, h=3.8, font_size=16):
    table = slide.shapes.add_table(
        len(rows), len(rows[0]), Inches(x), Inches(y), Inches(sum(widths)), Inches(h)
    ).table
    for i, width in enumerate(widths):
        table.columns[i].width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PALE_BLUE if r == 0 else WHITE
            cell.border = None
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Aptos"
            p.font.size = Pt(font_size)
            p.font.bold = r == 0
            p.font.color.rgb = NAVY if r == 0 else DARK
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
    return table


def new_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, title, subtitle)
    add_footer(slide, len(prs.slides))
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 — title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    title = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.15), Inches(11.8), Inches(1.2)
    )
    p = title.text_frame.paragraphs[0]
    p.text = "CDC-Aware RTL Generation Benchmark"
    p.font.name = "Aptos Display"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    sub = slide.shapes.add_textbox(
        Inches(0.82), Inches(2.35), Inches(11.6), Inches(0.7)
    )
    p = sub.text_frame.paragraphs[0]
    p.text = "Can an LLM generate RTL that is both functionally correct and CDC/RDC-safe?"
    p.font.name = "Aptos"
    p.font.size = Pt(23)
    p.font.color.rgb = DARK
    add_callout(
        slide,
        "Latest result: GPT-5.6 Sol had 29 functional passes; JasperGold passed 0 of them.",
        y=3.65,
    )
    who = slide.shapes.add_textbox(
        Inches(0.85), Inches(5.3), Inches(11.5), Inches(0.55)
    )
    p = who.text_frame.paragraphs[0]
    p.text = "Faiza · Advisor update · 24 August 2026"
    p.font.name = "Aptos"
    p.font.size = Pt(18)
    p.font.color.rgb = MUTED
    add_footer(slide, 1)

    # 2 — motivation
    slide = new_slide(prs, "Motivation and research question")
    add_bullets(
        slide,
        [
            "Most RTL-generation work scores compilation and functional simulation.",
            "Ordinary digital simulation does not model metastability or structural CDC/RDC safety.",
            "A design may pass its testbench while retaining unsafe clock/reset crossings.",
            "Research question: can generated RTL satisfy both function and CDC/RDC constraints?",
        ],
        size=22,
    )
    add_callout(
        slide,
        "Supported claim: functional pass does not imply CDC/RDC-clean.",
        y=5.75,
    )

    # 3 — benchmark
    slide = new_slide(prs, "Benchmark and evaluation pipeline")
    add_bullets(
        slide,
        [
            "44 circuits packaged with RTL, TB, SDC, and Jasper.",
            "Original and CDC/RDC-fixed RTL kept separate.",
            "Golden RTL is never in the prompt; generated RTL is never edited.",
            "Pilot: cdc_2phase, async_fifo, apbxclk.",
            "Functional prompt: ports and behavior.",
            "CDC-explicit: same, plus must be CDC/RDC-safe — no Gray/2-flop hints.",
        ],
        x=0.75,
        y=1.35,
        w=5.9,
        h=4.9,
        size=18,
    )
    flow = slide.shapes.add_textbox(
        Inches(6.9), Inches(1.5), Inches(5.5), Inches(4.7)
    )
    tf = flow.text_frame
    tf.clear()
    stages = [
        "Frozen prompt",
        "Model-generated RTL",
        "Xcelium compile",
        "Functional simulation",
        "JasperGold CDC/RDC",
        "Pass@k + archived evidence",
    ]
    for i, stage in enumerate(stages):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("↓  " if i else "") + stage
        p.font.name = "Aptos"
        p.font.size = Pt(21)
        p.font.bold = i in (0, 4, 5)
        p.font.color.rgb = RED if i == 4 else NAVY
        p.space_after = Pt(12)

    # 4 — expanded catalog
    slide = new_slide(
        prs,
        "Packaged CDC circuits: 44/44 tool-clean",
        "Same layout as the pilot: original RTL, fixed RTL, testbench, SDC, Jasper",
    )
    rows = [
        ["Set", "Count", "Functional sim", "JasperGold CDC/RDC"],
        ["Frozen pilot", "11", "11/11 PASS", "11/11 CLEAN"],
        ["Catalog imports", "19", "19/19 PASS", "19/19 CLEAN"],
        ["Additional CDC circuits", "14", "14/14 PASS", "14/14 CLEAN"],
        ["Total packaged", "44", "44/44 PASS", "44/44 CLEAN"],
    ]
    add_table(slide, rows, [3.4, 1.8, 3.2, 3.3], y=1.45, h=3.55, font_size=16)
    add_bullets(
        slide,
        [
            "Families added: Gray and 2-phase FIFOs, reset controllers, isochronous handshake, APB/AXI-Lite CDC, synchronizers.",
            "These 44 are packaged references, not 44 LLM results. Specs for the extra 33 remain draft.",
        ],
        y=5.15,
        h=1.55,
        size=17,
    )

    # 5 — one-shot results
    slide = new_slide(
        prs,
        "Frozen one-shot baselines",
        "Same scope: 3 circuits × 2 prompts × 3 attempts = 18 per model",
    )
    rows = [
        ["Model", "Base", "Compile", "Functional", "CDC-clean"],
        ["Llama 3.3 70B Instruct", "Llama", "4/18", "0/18", "0/18"],
        ["RTLCoder-v1.1 4-bit GGUF", "Mistral", "1/18", "0/18", "0/18"],
        ["RTLCoder-DeepSeek-v1.1 fp16", "DeepSeek", "5/18", "0/18", "0/18"],
    ]
    add_table(slide, rows, [4.4, 1.8, 1.8, 2.0, 1.8], y=1.5, h=2.7, font_size=16)
    add_bullets(
        slide,
        [
            "GGUF RTLCoder is Mistral-based. RTLCoder-DeepSeek is a separate fine-tune.",
            "Typical failures: invalid assignments, truncation, broken APB, CDC timeout.",
            "CDC-explicit wording alone did not produce a functional pass. Repair and Pass@20 are separate columns.",
        ],
        y=4.45,
        h=1.9,
        size=16,
    )

    # 6 — GPT result
    slide = new_slide(
        prs,
        "Main result: GPT-5.6 Sol exposes the evaluation gap",
        "3 circuits × 2 prompts × 10 attempts · compile → sim → Jasper after sim pass",
    )
    rows = [
        ["Circuit", "Prompt", "Compile", "Functional", "Jasper-clean"],
        ["cdc_2phase", "functional", "10/10", "10/10", "0/10"],
        ["cdc_2phase", "cdc_explicit", "10/10", "10/10", "0/10"],
        ["async_fifo", "functional", "10/10", "0/10", "not run"],
        ["async_fifo", "cdc_explicit", "10/10", "0/10", "not run"],
        ["apbxclk", "functional", "9/9", "9/9", "0/9"],
        ["apbxclk", "cdc_explicit", "—", "—", "credits exhausted"],
    ]
    add_table(slide, rows, [2.4, 2.4, 2.0, 2.3, 2.6], y=1.4, h=4.15, font_size=14)
    add_callout(
        slide,
        "29/29 sim passes failed Jasper. Typical tags: RST_PH_GLCH, RDC_RS_DFRS.",
        y=5.8,
    )

    # 7 — other protocols
    slide = new_slide(
        prs,
        "Other protocols (not mixed into the Sol table)",
        "One-shot, extra sampling, and repair stay in separate columns",
    )
    rows = [
        ["Protocol", "Result", "Takeaway"],
        [
            "Llama compile-repair",
            "6/6 compile; 1/6 functional; 0 clean",
            "Feedback fixes syntax, not CDC",
        ],
        [
            "RTLCoder-DeepSeek Pass@20",
            "73 RTL; 7 compile; 0 functional",
            "More samples found no sim pass",
        ],
        [
            "Jasper on those 7 compiles",
            "5 RDC fail; 2 zero crossings",
            "Zero violations ≠ automatically clean",
        ],
    ]
    add_table(slide, rows, [3.4, 4.4, 3.9], y=1.55, h=3.6, font_size=15)

    # 8 — limitations
    slide = new_slide(prs, "Limitations")
    add_bullets(
        slide,
        [
            "LLM generation is still the 3-circuit pilot, not all 44 packaged circuits.",
            "GPT-5.6 Sol is 49/60 complete; apbxclk CDC-explicit and one functional attempt remain.",
            "OpenAI GPT-5.6 does not expose temperature, top-p, or seed.",
            "RTLCoder-DeepSeek Pass@20 has no apbxclk yet.",
            "Specs for the extra 33 circuits are still draft.",
            "No synthesis, area, or assertion score yet.",
        ],
        size=20,
    )

    # 9 — future plan
    slide = new_slide(prs, "Proposed next phase")
    add_bullets(
        slide,
        [
            "Ask for this meeting",
            ("Resume the remaining 11 GPT-5.6 Sol attempts after credits are added.", 1),
            ("Run one more API model on the same 3 × 2 × 10 loop.", 1),
            "After that",
            ("Finish RTLCoder-DeepSeek Pass@20 when GPU is available.", 1),
            ("Expand generation from the 3-circuit pilot across the packaged 44.", 1),
        ],
        y=1.45,
        h=5.4,
        size=22,
    )

    # 11 — close
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    title = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.05), Inches(11.8), Inches(1.0)
    )
    p = title.text_frame.paragraphs[0]
    p.text = "Thank you"
    p.font.name = "Aptos Display"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    add_callout(slide, "Questions?", y=3.25)
    note = slide.shapes.add_textbox(
        Inches(1.1), Inches(4.35), Inches(11.1), Inches(1.1)
    )
    p = note.text_frame.paragraphs[0]
    p.text = "Compilation and simulation are necessary, but not sufficient, for generated CDC RTL."
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK
    who = slide.shapes.add_textbox(
        Inches(0.85), Inches(5.7), Inches(11.5), Inches(0.45)
    )
    p = who.text_frame.paragraphs[0]
    p.text = "Faiza · CDC-Aware RTL Generation · 24 August 2026"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(16)
    p.font.color.rgb = MUTED
    add_footer(slide, len(prs.slides))

    alt = ROOT / "slides" / "CDC_Aware_RTL_Generation.pptx"
    prs.save(OUT)
    prs.save(alt)
    print(f"Wrote {OUT}")
    print(f"Wrote {alt}")


if __name__ == "__main__":
    main()
